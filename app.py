import streamlit as st
import streamlit.components.v1 as components
import io, os, json, re, base64, uuid
from concurrent.futures import ThreadPoolExecutor

st.set_page_config(page_title="인사총무팀 AI 에이전트", page_icon="🤖", layout="wide",
                   initial_sidebar_state="expanded")

CFG_FILE = "wellfine_config.json"
SHARED_HISTORY_FILE = "shared_history.json"

def load_cfg():
    try:
        if os.path.exists(CFG_FILE):
            with open(CFG_FILE, encoding="utf-8") as f:
                return json.load(f)
    except Exception:
        pass
    return {}

def save_cfg(data):
    try:
        with open(CFG_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False)
    except Exception:
        pass

def load_shared_history():
    try:
        if os.path.exists(SHARED_HISTORY_FILE):
            with open(SHARED_HISTORY_FILE, encoding="utf-8") as f:
                return json.load(f)
    except Exception:
        pass
    return []

def save_shared_history(history):
    try:
        tmp_file = SHARED_HISTORY_FILE + ".tmp"
        with open(tmp_file, "w", encoding="utf-8") as f:
            json.dump(history, f, ensure_ascii=False)
        os.replace(tmp_file, SHARED_HISTORY_FILE)
    except Exception:
        pass

def _entry_key(entry):
    return entry.get("id") or "|".join([
        entry.get("time", ""), entry.get("title", ""),
        entry.get("doc", ""), entry.get("question", "")
    ])

def append_shared_history(entry):
    entry = dict(entry)
    entry.setdefault("id", uuid.uuid4().hex)
    latest = load_shared_history()
    keys = {_entry_key(h) for h in latest}
    if _entry_key(entry) not in keys:
        latest.append(entry)
        save_shared_history(latest)
    st.session_state.meeting_history = latest
    return max(0, len(latest) - 1)

def sync_shared_history():
    if st.session_state.get("meeting_running"):
        return
    latest = load_shared_history()
    if latest != st.session_state.get("meeting_history", []):
        st.session_state.meeting_history = latest

st.markdown("""
<style>
[data-testid="stAppViewContainer"]{background:#f5f1eb;}
[data-testid="stHeader"]{background:transparent;}
[data-testid="stSidebar"]{background:#f3efe8; }
[data-testid="stSidebarContent"]{padding:14px 10px 16px;}
#MainMenu,footer{visibility:hidden;}
.sidebar-brand{display:flex;align-items:center;gap:8px;margin:2px 4px 12px;
  font-size:18px;font-weight:800;color:#211c18;}
.sidebar-brand .mark{width:24px;height:24px;border-radius:8px;background:#26211c;
  color:#fff;display:inline-flex;align-items:center;justify-content:center;font-size:14px;}
.sidebar-section{font-size:11px;font-weight:800;color:#877d72;margin:18px 8px 7px;
  text-transform:uppercase;letter-spacing:.04em;}
.history-empty{font-size:12px;line-height:1.55;color:#8a8178;background:#fffaf2;
  border:1px dashed #d8cfc3;border-radius:10px;padding:12px;margin:8px 2px;}
[data-testid="stSidebar"] div[data-testid="stButton"] > button{
  min-height:34px;border-radius:10px;border:1px solid #ded4c7;background:#fbf7ef;
  color:#2a2520;font-weight:650;justify-content:flex-start;text-align:left;
  white-space:normal;}
[data-testid="stSidebar"] div[data-testid="stButton"] > button:hover{
  border-color:#b9aa99;background:#eee6db;color:#111;}
[data-testid="stSidebar"] div[data-testid="stTextInput"] input{
  border-radius:10px;background:#fbf7ef;border-color:#ded4c7;}
</style>
""", unsafe_allow_html=True)

PROVIDERS = {
    "🟠 Claude (Anthropic)": {"key": "anthropic", "models": {
        "Claude Sonnet (기본)": "claude-sonnet-4-6",
        "Claude Opus (고성능)": "claude-opus-4-8",
        "Claude Haiku (빠름)": "claude-haiku-4-5-20251001",
    }},
    "🟢 ChatGPT (OpenAI)": {"key": "openai", "models": {
        "GPT-4o (권장)": "gpt-4o",
        "GPT-4o mini (빠름)": "gpt-4o-mini",
    }},
    "🔵 Gemini (Google)": {"key": "gemini", "models": {
        "Gemini 2.5 Flash (무료/빠름)": "gemini-2.5-flash",
        "Gemini 2.5 Pro (고성능)": "gemini-2.5-pro",
        "Gemini 1.5 Flash": "gemini-1.5-flash",
    }},
}

BOTS = [
    {"name": "노무봇",  "system": "웰파인(강원도 횡성 건강기능식품 OEM/ODM) 인사총무팀 노무 전문 AI.\n성격: 보수적이고 까다로운 노무사 스타일.\n- 법 조문 근거 없이 괴다고 하지 마세요\n- 애매하면 반드시 전문가 확인 필요 명시\n- 에 노동관련 법령 관점으로 검토\n- 핵심 포인트 번호 목록으로 제시"},
    {"name": "HRM봇",  "system": "웰파인 인사총무팀 HRM 전문 AI.\n성격: 원칙주의자.\n- 사내 규정, 취업규칙과의 일치 여부 체크\n- 규정과 실제 운영 간의 괴리를 찾아내세요\n- 형평성 체크\n- 핵심 포인트 번호 목록으로 제시"},
    {"name": "채용봇",  "system": "웰파인 인사총무팀 채용 전문 AI.\n성격: 현실주의자.\n- 강원도 횡성 지역의 지원자 풀이 제한적임 감안\n- JD의 모호한/차별적 표현 찾아내세요\n- 채용 과정의 법적 문제 체크\n- 핵심 포인트 번호 목록으로 제시"},
    {"name": "계약봇",  "system": "웰파인 인사총무팀 계약 전문 AI.\n성격: 의심 많은 변호사 스타일.\n- 계약서 조항 꼼꼼하게 분석\n- 불리한 조항/모호한 표현 찾아내세요\n- 해지/갱신 조항 리스크 집중 체크\n- 핵심 포인트 번호 목록으로 제시"},
    {"name": "총무봇",  "system": "웰파인 인사총무팀 총무 전문 AI.\n성격: 꼼꼼한 살림꼼.\n- 비용 처리, 시설, 물품 관련 규정 준수 체크\n- 세무적으로 문제없는지 검토\n- 내부 결재/승인 프로세스 확인\n- 핵심 포인트 번호 목록으로 제시"},
    {"name": "급여봇",  "system": "웰파인 인사총무팀 급여/4대보험 전문 AI.\n성격: 숫자에 목숙 거는 스타일.\n- 급여 계산과 공제 항목 정확성 체크\n- 4대보험 요율과 신고 기한 확인\n- E-9 비자 외국인 근로자의 특수한 처리 체크\n- 핵심 포인트 번호 목록으로 제시"},
    {"name": "교육봇",  "system": "웰파인 인사총무팀 교육/OJT 전문 AI.\n성격: 실용주의 교육자.\n- 교육 내용이 실제 업무와 연결되는지 체크\n- OJT/수습평가 설계의 공정성 확인\n- 외국인 근로자 교육 시 언어/문화 차이 고려\n- 핵심 포인트 번호 목록으로 제시"},
]

FACTCHECK_SYS = (
    "웰파인 인사총무팀 팩트체크 AI. 감정 없는 검사 스타일.\n\n"
    "반드시 아래 실제 법령/규정 기준으로 검증하세요:\n"
    "- 근로기준법\n- 최저임금법\n- 근로자퇴직급여보장법\n"
    "- 고용보험법 / 산재보험법 / 국민연금법 / 국민건강보험법\n"
    "- 여성고용평등법\n- 외국인근로자 고용법 (E-9 비자)\n"
    "- 채용절차 공정화에 관한 법률\n- 개인정보 보호법\n\n"
    "각 봇 의견에 대해 [맞음/틀림/확인필요] 판정 후 근거 법령 조문 명시."
)

FIELD_SYS = (
    "웰파인 인사총무팀 실무 전문 AI. 성격: 현장 경험 많은 선배.\n"
    "- '이론은 맞는데 실제로 가능해?' 관점으로 접근\n"
    "- 직원 입장에서 이 결정을 받아들일 수 있는지 시뮬레이션\n"
    "- 강원도 횡성 제조업 현장 특성 감안\n"
    "- 현실적으로 실행 가능한 방향 제시"
)

FINAL_SYS = (
    "웰파인 인사총무팀 AI 회의 최종 정리 AI. 냉철한 회의 진행자.\n"
    "- 모든 봇 의견 종합해 결론 도출\n"
    "- 의견 충돌한 부분은 양쪽 모두 명시\n"
    "- 구체적 액션 아이템을 우선순위 순서로 제시\n"
    "- 전문가 확인 반드시 필요한 사항 명시\n\n"
    "다음 형식으로 작성:\n"
    "## 핵심 결론\n## 주요 리스크\n## 액션 아이템 (우선순위 순)\n## 전문가 확인 필요 사항"
)

SCREEN_SYS = (
    "웰파인 HR AI 회의 코디네이터. 이번 안건과 관련 있는 봇만 골라주세요.\n\n"
    "봇 전문 영역:\n"
    "- 노무봇: 근로기준법, 수습/해고/징계, 근로시간, 퇴직, 실업급여, 노동 법령\n"
    "- HRM봇: 인사관리, 사내규정, 취업규칙, 인사제도, 평가\n"
    "- 채용봇: 채용공고/JD, 면접, 지원자 선발, 채용 프로세스\n"
    "- 계약봇: 근로계약서, 용역계약, 협약서 등 계약 문서\n"
    "- 총무봇: 비용처리/법인카드, 시설관리, 물품/소모품, 세금계산서\n"
    "- 급여봇: 급여 계산, 임금, 4대보험, 원천징수, 퇴직금, 수당\n"
    "- 교육봇: OJT, 수습평가, 직원교육, 교육훈련\n\n"
    '{"\ub178\ubb34\ubd07":true,"HRM\ubd07":false,"\ucc44\uc6a9\ubd07":false,"\uacc4\uc57d\ubd07":false,'
    '"\ucd1d\ubb34\ubd07":false,"\uae09\uc5ec\ubd07":false,"\uad50\uc721\ubd07":false} \ud615\ud0dc\ub85c\ub9cc \uc751\ub2f5\ud558\uc138\uc694.'
)

MEMORY_SYS = (
    "웰파인 HR 회의 요약 AI. "
    "이번 회의에서 향후 참고할 핵심 사항 2-3줄로만 요약. "
    "날짜 불필요. 간결하게."
)

BOT_COLORS = {
    "\ub178\ubb34\ubd07":     ["#3B6FD4","#2756B8"],
    "HRM\ubd07":      ["#2E9E55","#1E7D3E"],
    "\ucc44\uc6a9\ubd07":     ["#E07B00","#B86000"],
    "\uacc4\uc57d\ubd07":     ["#7048E8","#5030C8"],
    "\ucd1d\ubb34\ubd07":     ["#546E8A","#3D5570"],
    "\uae09\uc5ec\ubd07":     ["#C79800","#A07800"],
    "\uad50\uc721\ubd07":     ["#0A9E7A","#087D60"],
    "\ud329\ud2b8\uccb4\ud06c\ubd07": ["#C92A2A","#A01E1E"],
    "\uc2e4\ubb34\ubd07":     ["#7B5E3C","#5A4228"],
    "\ucd5c\uc885\uc815\ub9ac\ubd07": ["#1565C0","#0D47A0"],
}

def call_bot(pkey, api_key, model, system, content, max_tokens=700):
    combined = system + "\n\n" + content
    try:
        if pkey == "anthropic":
            import anthropic
            client = anthropic.Anthropic(api_key=api_key)
            resp = client.messages.create(model=model, max_tokens=max_tokens,
                system=system, messages=[{"role": "user", "content": content}])
            return resp.content[0].text
        elif pkey == "openai":
            from openai import OpenAI
            client = OpenAI(api_key=api_key)
            resp = client.chat.completions.create(model=model, max_tokens=max_tokens,
                messages=[{"role": "system", "content": system},
                          {"role": "user", "content": content}])
            return resp.choices[0].message.content
        elif pkey == "gemini":
            import google.generativeai as genai
            genai.configure(api_key=api_key)
            resp = genai.GenerativeModel(model).generate_content(combined)
            return resp.text
    except Exception as e:
        return f"오류: {e}"
    return "지원하지 않는 제공자."

def extract_text(f):
    if not f:
        return ""
    ext = f.name.split(".")[-1].lower()
    try:
        if ext == "txt":
            return f.read().decode("utf-8", errors="ignore")
        elif ext == "pdf":
            import pypdf
            r = pypdf.PdfReader(io.BytesIO(f.read()))
            return "\n".join(p.extract_text() or "" for p in r.pages)
        elif ext in ("docx", "doc"):
            from docx import Document
            d = Document(io.BytesIO(f.read()))
            return "\n".join(p.text for p in d.paragraphs)
        elif ext in ("pptx", "ppt"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(f.read()))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
            return "\n".join(texts)
    except Exception as e:
        return f"[파일 오류: {e}]"
    return ""

def build_content(doc_text, question, case_notes):
    parts = []
    if case_notes.strip():
        parts.append("[웰파인 누적 케이스]\n" + case_notes.strip())
    if doc_text.strip():
        parts.append("[문서]\n" + doc_text[:3000])
    if question.strip():
        parts.append("[질문/검토 요청]\n" + question.strip())
    return "\n\n".join(parts)

def build_partial_summary(results):
    if not results:
        return ""
    labels = {
        "팩트체크봇_1차": "팩트체크봇 1차",
        "팩트체크봇_2차": "팩트체크봇 최종",
    }
    lines = [
        "## 중단 시점까지의 검토 내용",
        "아래 내용은 회의가 중단되기 전까지 완료된 봇 발언입니다.",
    ]
    for name, result in results.items():
        if not result:
            continue
        title = labels.get(name, name)
        lines.append(f"\n### {title}\n{result}")
    return "\n".join(lines).strip()


def render_room(statuses, results, doc_name=""):
    SC = {'idle':'#bbb','active':'#f5c842','done':'#52c463','error':'#e05555'}

    # 말풍선 모달용 발언 데이터 (results 그대로 사용, 팩트체크봇은 최신 키로 매핑)
    _speech = dict(results or {})
    if '팩트체크봇_2차' in _speech:
        _speech['팩트체크봇'] = _speech['팩트체크봇_2차']
    elif '팩트체크봇_1차' in _speech:
        _speech['팩트체크봇'] = _speech['팩트체크봇_1차']
    _speech_json = json.dumps(_speech, ensure_ascii=False)

    def robot_face(name, C, D):
        if '노무' in name:
            return f'<path d="M21 31l8 2.4M43 31l-8 2.4" stroke="#273445" stroke-width="2.1" stroke-linecap="round"/><circle cx="25" cy="37" r="3.6" fill="{C}"/><circle cx="39" cy="37" r="3.6" fill="{C}"/><path d="M28.5 44.5h7" stroke="#273445" stroke-width="2.1" stroke-linecap="round"/>'
        if '계약' in name:
            return f'<path d="M21.5 31.5l7.2-.8M35.2 31.2l8 2.1" stroke="#273445" stroke-width="2" stroke-linecap="round"/><circle cx="25" cy="37" r="3.7" fill="{C}"/><circle cx="39" cy="37" r="3.1" fill="{C}"/><path d="M29 44c2.6 1.5 6.2 1.1 8-.7" fill="none" stroke="#273445" stroke-width="2" stroke-linecap="round"/>'
        if '팩트' in name:
            return f'<rect x="20.5" y="33" width="9" height="6.5" rx="3.2" fill="{C}"/><rect x="34.5" y="33" width="9" height="6.5" rx="3.2" fill="{C}"/><path d="M29.5 36.2h5" stroke="{D}" stroke-width="1.8" stroke-linecap="round"/><path d="M28.5 44h7" stroke="#273445" stroke-width="2.1" stroke-linecap="round"/>'
        if '급여' in name:
            return f'<circle cx="25" cy="36" r="3.9" fill="{C}"/><circle cx="39" cy="36" r="3.9" fill="{C}"/><circle cx="26.4" cy="34.6" r="1.05" fill="#FFFFFF"/><circle cx="40.4" cy="34.6" r="1.05" fill="#FFFFFF"/><path d="M29 43.5h6" stroke="#273445" stroke-width="2" stroke-linecap="round"/>'
        if 'HRM' in name:
            return f'<path d="M21 32.5l9 1.2M43 32.5l-9 1.2" stroke="#273445" stroke-width="2.1" stroke-linecap="round"/><rect x="22" y="36" width="7" height="5" rx="3.5" fill="{C}"/><rect x="35" y="36" width="7" height="5" rx="3.5" fill="{C}"/><path d="M29 44h6" stroke="#273445" stroke-width="2" stroke-linecap="round"/>'
        if '채용' in name:
            return f'<circle cx="25" cy="36.5" r="3.8" fill="{C}"/><circle cx="39" cy="36.5" r="3.8" fill="{C}"/><path d="M29 43c1.8 2.2 5.5 2.2 6 0" fill="none" stroke="#273445" stroke-width="2" stroke-linecap="round"/>'
        if '총무' in name:
            return f'<path d="M22 33.5l8.5 2M42 33.5l-8.5 2" stroke="#273445" stroke-width="2.1" stroke-linecap="round"/><circle cx="25.5" cy="38.5" r="3.5" fill="{C}"/><circle cx="38.5" cy="38.5" r="3.5" fill="{C}"/><path d="M29 44h6" stroke="#273445" stroke-width="2" stroke-linecap="round"/>'
        if '교육' in name:
            return f'<circle cx="25" cy="36" r="3.8" fill="{C}"/><circle cx="39" cy="36" r="3.8" fill="{C}"/><path d="M27.5 44.5h9" stroke="#273445" stroke-width="2.5" stroke-linecap="round"/>'
        if '실무' in name:
            return f'<path d="M21 32l9 3M43 32l-9 3" stroke="#273445" stroke-width="2.1" stroke-linecap="round"/><circle cx="25" cy="38" r="3.7" fill="{C}"/><circle cx="39" cy="38" r="3.7" fill="{C}"/><path d="M28.5 44c1 1.8 5 1.8 7 0" fill="none" stroke="#273445" stroke-width="2" stroke-linecap="round"/>'
        return f'<circle cx="25" cy="36" r="3.8" fill="{C}"/><circle cx="39" cy="36" r="3.8" fill="{C}"/><path d="M29 44h6" stroke="#273445" stroke-width="2" stroke-linecap="round"/>'

    def robot_accessory(name, C, D):
        if '노무' in name:
            return f'<rect x="17" y="18" width="30" height="9" rx="4.5" fill="#273445"/><rect x="19" y="20" width="26" height="5" rx="2.5" fill="{C}" opacity=".35"/><circle cx="32" cy="14" r="3.5" fill="#273445"/>'
        if '계약' in name:
            return f'<path d="M17 23a15 5 0 0 1 30 0" fill="{C}"/><rect x="16" y="21.5" width="32" height="4" rx="2" fill="#273445"/>'
        if '팩트' in name:
            return f'<rect x="19" y="15" width="26" height="11" rx="3" fill="#273445"/><rect x="21" y="17.5" width="10" height="2" rx="1" fill="{C}"/><rect x="21" y="21" width="16" height="2" rx="1" fill="{C}" opacity=".5"/>'
        if '급여' in name:
            return f'<rect x="25" y="12" width="14" height="14" rx="3" fill="#273445"/><text x="32" y="23" font-size="9" fill="{C}" text-anchor="middle" font-weight="bold">₩</text>'
        if 'HRM' in name:
            return f'<rect x="26" y="13" width="12" height="10" rx="5" fill="#273445"/><rect x="28" y="16" width="8" height="4" rx="2" fill="{C}" opacity=".6"/>'
        if '채용' in name:
            return f'<path d="M22 23l10-9 10 9" stroke="#273445" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round" fill="none"/><rect x="26" y="16" width="12" height="9" rx="2" fill="{C}" opacity=".4"/>'
        if '총무' in name:
            return f'<circle cx="32" cy="17" r="7" fill="#273445"/><circle cx="32" cy="17" r="4" fill="{C}" opacity=".5"/>'
        if '교육' in name:
            return f'<rect x="20" y="14" width="24" height="12" rx="3" fill="#273445"/><path d="M24 18h8M24 22h12" stroke="{C}" stroke-width="1.5" stroke-linecap="round"/>'
        if '실무' in name:
            return f'<path d="M20 22 Q32 13 44 22" fill="{D}"/><rect x="19" y="21" width="26" height="4" rx="2" fill="#273445"/>'
        return f'<polygon points="32,10 44,20 40,22 32,14 24,22 20,20" fill="#273445"/><polygon points="32,10 44,20 40,22 32,14 24,22 20,20" fill="{C}" opacity=".4"/>'

    def svg_robot(name, sz):
        cols = BOT_COLORS.get(name, ['#888888','#666666'])
        C, D = cols
        return (f'<svg width="{sz}" height="{sz}" viewBox="0 0 64 70" xmlns="http://www.w3.org/2000/svg" style="filter:drop-shadow(0 4px 7px rgba(36,28,20,.22));overflow:visible">'
                f'<ellipse cx="32" cy="63" rx="19" ry="4" fill="#2d241b" opacity=".16"/>'
                f'<rect x="13" y="24" width="38" height="32" rx="10" fill="{D}"/>'
                f'<rect x="15" y="26" width="34" height="28" rx="9" fill="#F4EFE6"/>'
                + robot_face(name, C, D)
                + robot_accessory(name, C, D)
                + f'<rect x="22" y="54" width="20" height="7" rx="3.5" fill="{D}" opacity=".92"/>'
                f'<circle cx="32" cy="57.5" r="2.3" fill="#fff" opacity=".75"/>'
                f'<path d="M20 57v7" stroke="#9B8E7D" stroke-width="3" stroke-linecap="round"/>'
                f'<path d="M44 57v7" stroke="#9B8E7D" stroke-width="3" stroke-linecap="round"/>'
                f'<path d="M8 38h5M51 38h5" stroke="#9B8E7D" stroke-width="3.2" stroke-linecap="round"/>'
                f'</svg>')

    def bot_card(name, status):
        color = SC.get(status, '#bbb')
        glow = f'box-shadow:0 0 8px {color};' if status != 'idle' else ''
        bubble = (f'<div onclick="wfShowSpeech(\'{name}\')" '
                  'style="cursor:pointer;font-size:9px;font-weight:700;background:#fff8e1;'
                  'border:1px solid #f0d98a;border-radius:9px;padding:2px 7px;color:#8a6500;'
                  'margin-bottom:2px;white-space:nowrap;box-shadow:0 1px 3px rgba(0,0,0,.12);'
                  'animation:wfBob 1.1s ease-in-out infinite;">\U0001F4AC 발언중</div>') if status == 'active' else ''
        short = name.replace('제크', '').replace('정리', '')
        return (f'<div style="display:flex;flex-direction:column;align-items:center;gap:2px;padding:2px 5px;min-width:66px;">'
                + bubble
                + svg_robot(name, 62)
                + f'<div style="font-size:10px;font-weight:800;color:#2a2520;text-align:center;line-height:1.1">{short}</div>'
                + f'<div style="width:7px;height:7px;border-radius:50%;background:{color};{glow}"></div>'
                + '</div>')

    def g(n): return statuses.get(n, 'idle')

    top_row = ''.join(bot_card(n, g(n)) for n in ['노무봇','HRM봇','채용봇','계약봇'])
    btm_row = ''.join(bot_card(n, g(n)) for n in ['총무봇','급여봇','교육봇'])

    if doc_name:
        table_inner = f'<div style="background:rgba(255,255,255,.88);border-radius:4px;padding:3px 10px;font-size:9px;color:#3a2800;max-width:220px;white-space:nowrap;overflow:hidden;text-overflow:ellipsis">📄 {doc_name}</div>'
    else:
        table_inner = '<div style="height:46px"></div>'

    cs = g('최종정리봇')
    c_color = SC.get(cs, '#bbb')
    c_glow = f'box-shadow:0 0 7px {c_color};' if cs != 'idle' else ''
    c_bubble = ('<div onclick="wfShowSpeech(\'최종정리봇\')" '
                'style="cursor:pointer;font-size:9px;font-weight:700;background:#fff8e1;'
                'border:1px solid #f0d98a;border-radius:9px;padding:2px 7px;color:#8a6500;'
                'margin-bottom:2px;white-space:nowrap;box-shadow:0 1px 3px rgba(0,0,0,.12);'
                'animation:wfBob 1.1s ease-in-out infinite;">\U0001F4AC 발언중</div>') if cs == 'active' else ''

    podium = (
        '<div style="display:flex;flex-direction:column;align-items:center;justify-content:center;'
        'min-width:124px;margin-left:8px;">'
        '<div style="display:flex;flex-direction:column;align-items:center;gap:2px;'
        'background:rgba(255,255,255,.48);border:1px solid #d7cfc3;border-radius:12px;'
        'padding:8px 12px 10px;box-shadow:0 5px 14px rgba(44,35,24,.12);">'
        + c_bubble
        + svg_robot('최종정리봇', 72)
        + '<div style="width:72px;height:24px;margin-top:-6px;border-radius:7px 7px 4px 4px;'
          'background:linear-gradient(180deg,#7C5838,#50331F);border:2px solid #3d2618;'
          'box-shadow:0 5px 10px rgba(36,25,16,.25);display:flex;align-items:center;'
          'justify-content:center;color:rgba(255,255,255,.62);font-size:8px;font-weight:800;'
          'letter-spacing:1px;">HOST</div>'
        + '<div style="font-size:10px;font-weight:800;color:#2a2520;margin-top:3px">최종정리봇</div>'
          '<div style="font-size:9px;color:#6b5e4f">회의 진행자</div>'
        + f'<div style="width:8px;height:8px;border-radius:50%;background:{c_color};{c_glow}"></div>'
        + '</div></div>'
    )

    # 말풍선 클릭 시 뜨는 모달 + 애니메이션 (iframe 내부 자체 완결)
    _modal_html = (
        '<style>@keyframes wfBob{0%,100%{transform:translateY(0)}50%{transform:translateY(-3px)}}</style>'
        '<div id="wfModal" onclick="if(event.target===this)this.style.display=\'none\'" '
        'style="display:none;position:fixed;inset:0;background:rgba(25,18,10,.5);z-index:99999;'
        'align-items:center;justify-content:center;padding:14px;box-sizing:border-box;'
        'font-family:-apple-system,sans-serif;">'
        '<div style="background:#fffdf7;max-width:460px;width:100%;max-height:86%;overflow:auto;'
        'border-radius:14px;padding:16px 18px;box-shadow:0 14px 44px rgba(0,0,0,.4);border:1px solid #e6dcc8;">'
        '<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:10px;">'
        '<div id="wfModalTitle" style="font-size:14px;font-weight:800;color:#2a2520;"></div>'
        '<div onclick="document.getElementById(\'wfModal\').style.display=\'none\'" '
        'style="cursor:pointer;font-size:20px;color:#9a8b76;line-height:1;padding:0 4px;">×</div>'
        '</div>'
        '<div id="wfModalBody" style="font-size:12.5px;line-height:1.65;color:#3a342c;'
        'white-space:pre-wrap;word-break:break-word;"></div>'
        '</div></div>'
        '<script>'
        'var WF_SPEECH=' + _speech_json + ';'
        'function wfShowSpeech(n){'
        'document.getElementById("wfModalTitle").textContent="\U0001F4AC "+n+" 발언";'
        'var t=WF_SPEECH[n];'
        'document.getElementById("wfModalBody").textContent=(t&&(""+t).trim())?t:"아직 이 봇의 발언 내용이 도착하지 않았습니다.";'
        'document.getElementById("wfModal").style.display="flex";'
        '}'
        '</script>'
    )

    return (
        '<div style="background:linear-gradient(180deg,#e8e2d6,#dfd9cd);border-radius:16px;'
        'padding:18px 10px 16px;border:1px solid #cec8be;font-family:-apple-system,sans-serif;">'
        '<div style="display:flex;justify-content:center;align-items:center;gap:16px;">'
        '<div style="display:flex;flex-direction:column;align-items:center;">'
        '<div style="display:flex;justify-content:center;align-items:flex-end;gap:8px;margin-bottom:6px;">'
        + top_row + '</div>'
        '<div style="display:flex;justify-content:center;align-items:center;gap:12px;margin:4px 0;">'
        + bot_card('팩트체크봇', g('팩트체크봇'))
        + '<div style="width:320px;min-height:92px;background:linear-gradient(160deg,#6B4423,#8B5E3C 40%,#9B6B45 60%,#6B4423);'
          'border-radius:10px;border:3px solid #4A2E18;box-shadow:0 6px 20px rgba(0,0,0,.4);'
          'display:flex;flex-direction:column;align-items:center;justify-content:center;gap:4px;padding:8px;">'
        + table_inner + '</div>'
        + bot_card('실무봇', g('실무봇'))
        + '</div>'
        '<div style="display:flex;justify-content:center;align-items:flex-end;gap:8px;margin-top:6px;">'
        + btm_row + '</div>'
        + '</div>'
        + podium
        + '</div></div>'
        + _modal_html
    )


# ── 세션 초기화 ────────────────────────────────────────────
_cfg = load_cfg()
for k, v in [
    ("bot_statuses", {}), ("bot_results", {}), ("case_notes", ""),
    ("stop_requested", False), ("meeting_running", False),
    ("meeting_history", []), ("renaming_idx", None), ("current_doc", ""),
    ("show_result_panel", False), ("last_final", ""), ("selected_history_idx", None),
    ("current_question", ""), ("feedback_pending", ""), ("stop_notice", ""),
    ("saved_provider", _cfg.get("provider", list(PROVIDERS.keys())[0])),
    ("saved_model",    _cfg.get("model", "")),
    ("saved_api_key",  _cfg.get("api_key", "")),
]:
    if k not in st.session_state:
        st.session_state[k] = v

# 공유 히스토리: 세션 첫 로드 시 파일에서 불러옴
if "history_loaded" not in st.session_state:
    st.session_state.history_loaded = True
    _shared = load_shared_history()
    if _shared:
        st.session_state.meeting_history = _shared

if "bot_prompts" not in st.session_state:
    bp0 = {b["name"]: b["system"] for b in BOTS}
    bp0.update({"팩트체크봇": FACTCHECK_SYS, "실무봇": FIELD_SYS, "최종정리봇": FINAL_SYS})
    st.session_state.bot_prompts = bp0

# ── 현재 설정값 도출 ───────────────────────────────────────
_p = st.session_state.saved_provider
if _p not in PROVIDERS:
    _p = list(PROVIDERS.keys())[0]
_prov    = PROVIDERS[_p]
_pkey    = _prov["key"]
_ml      = st.session_state.saved_model
if _ml not in _prov["models"]:
    _ml = list(_prov["models"].keys())[0]
_model   = _prov["models"][_ml]
_api_key = st.session_state.saved_api_key
if not _api_key:
    try:
        smap = {"anthropic": "ANTHROPIC_API_KEY", "openai": "OPENAI_API_KEY", "gemini": "GEMINI_API_KEY"}
        _api_key = st.secrets.get(smap[_pkey], "")
    except Exception:
        _api_key = ""

@st.dialog("⚙️ 설정", width="large")
def show_settings():
    tab1, tab2, tab3 = st.tabs(["🔌 AI 연결", "🤖 봇 프롬프트", "📂 케이스 메모"])
    with tab1:
        pl = st.selectbox("AI 제공자", list(PROVIDERS.keys()),
            index=(list(PROVIDERS.keys()).index(st.session_state.saved_provider)
                   if st.session_state.saved_provider in PROVIDERS else 0))
        pt = PROVIDERS[pl]; pkt = pt["key"]
        mll = list(pt["models"].keys())
        saved_m = st.session_state.saved_model
        ml = st.selectbox("모델", mll, index=mll.index(saved_m) if saved_m in mll else 0)
        ph = {"anthropic": "sk-ant-...", "openai": "sk-...", "gemini": "AIza..."}.get(pkt, "")
        ak = st.text_input("API 키 (팀 공유)", type="password",
                           value=st.session_state.saved_api_key, placeholder=ph)
        st.caption("팀원 모두가 이 키로 사용합니다.")
        if st.button("💾 저장", type="primary", use_container_width=True):
            st.session_state.saved_provider = pl
            st.session_state.saved_model    = ml
            st.session_state.saved_api_key  = ak
            save_cfg({"provider": pl, "model": ml, "api_key": ak})
            st.rerun()
    with tab2:
        all_names = ([b["name"] for b in BOTS] + ["팩트체크봇", "실무봇", "최종정리봇"])
        bs = st.selectbox("봇 선택", all_names)
        ed = st.text_area("프롬프트", value=st.session_state.bot_prompts.get(bs, ""), height=220)
        if st.button("💾 프롬프트 저장", use_container_width=True):
            st.session_state.bot_prompts[bs] = ed
            st.rerun()
    with tab3:
        st.caption("모든 회의에 기본 콘텍스트로 포함됩니다.")
        nn = st.text_area("누적 케이스", value=st.session_state.case_notes, height=200,
                          placeholder="수습 불합격 시 면담 2회 이상 기록 필요...")
        if st.button("💾 케이스 저장", use_container_width=True):
            st.session_state.case_notes = nn
            st.rerun()

def _short(text, limit=28):
    text = (text or "").replace("\n", " ").strip()
    return text if len(text) <= limit else text[:limit - 1] + "…"

def _open_history(real_idx):
    h = st.session_state.meeting_history[real_idx]
    statuses = {}; results = {}
    if h.get("results") or h.get("statuses"):
        results = dict(h.get("results", {}))
        statuses = dict(h.get("statuses", {}))
    else:
        for name, result in h.get("ops2", {}).items():
            statuses[name] = "done"; results[name] = result
        if h.get("fc2"):
            statuses["팩트체크봇"] = "done"; results["팩트체크봇"] = h["fc2"]
        if h.get("final"):
            statuses["최종정리봇"] = "done"; results["최종정리봇"] = h["final"]
    panel_text = h.get("final") or h.get("partial_summary") or (
        build_partial_summary(results) if h.get("partial") else ""
    )
    st.session_state.bot_statuses = statuses
    st.session_state.bot_results = results
    st.session_state.current_doc = h.get("doc", "")
    st.session_state.last_final = panel_text
    st.session_state.show_result_panel = bool(panel_text)
    st.session_state.selected_history_idx = real_idx
    st.session_state.meeting_running = False
    st.session_state.stop_requested = False

def _new_meeting():
    st.session_state.bot_statuses = {}
    st.session_state.bot_results = {}
    st.session_state.current_doc = ""
    st.session_state.last_final = ""
    st.session_state.show_result_panel = False
    st.session_state.selected_history_idx = None
    st.session_state.renaming_idx = None
    st.session_state.meeting_running = False
    st.session_state.stop_requested = False
    st.session_state.feedback_pending = ""
    st.session_state.stop_notice = ""

with st.sidebar:
    sync_shared_history()
    st.markdown(
        '<div class="sidebar-brand"><span class="mark">HR</span><span>인사총무팀 AI 에이전트</span></div>',
        unsafe_allow_html=True)
    if st.button("＋ 새 회의", key="new_meeting", use_container_width=True):
        _new_meeting(); st.rerun()
    query = st.text_input("회의 검색", key="history_filter", placeholder="회의 검색",
                          label_visibility="collapsed")
    st.markdown('<div class="sidebar-section">최근 회의</div>', unsafe_allow_html=True)
    history = st.session_state.meeting_history
    if not history:
        st.markdown('<div class="history-empty">회의를 진행하면 이곳에 저장됩니다.</div>',
                    unsafe_allow_html=True)
    else:
        q = query.strip().lower()
        visible = []
        for real_idx, h in enumerate(history):
            haystack = " ".join([h.get("title",""), h.get("question",""),
                                 h.get("doc",""), h.get("final","")]).lower()
            if not q or q in haystack:
                visible.append(real_idx)
        if not visible:
            st.markdown('<div class="history-empty">검색 결과가 없습니다.</div>',
                        unsafe_allow_html=True)
        for real_idx in reversed(visible):
            h = history[real_idx]
            title = h.get("title") or h.get("question") or "새 회의"
            if h.get("partial"):
                title = "⏸ " + title
            meta_bits = [h.get("time","")]
            if h.get("doc"):
                meta_bits.append(h["doc"])
            elif h.get("question"):
                meta_bits.append(h["question"])
            meta = " · ".join(_short(x,22) for x in meta_bits if x)
            active = st.session_state.selected_history_idx == real_idx
            row_title = ("● " if active else "") + _short(title, 28)
            row_label = row_title + (f"\n{_short(meta,34)}" if meta else "")
            oc, rc, dc = st.columns([5.8, 1.1, 1.1])
            with oc:
                if st.button(row_label, key=f"hist_open_{real_idx}", use_container_width=True):
                    _open_history(real_idx); st.rerun()
            with rc:
                if st.button("✎", key=f"hist_ren_{real_idx}", use_container_width=True):
                    st.session_state.renaming_idx = real_idx; st.rerun()
            with dc:
                if st.button("×", key=f"hist_del_{real_idx}", use_container_width=True):
                    st.session_state.meeting_history.pop(real_idx)
                    save_shared_history(st.session_state.meeting_history)
                    if st.session_state.selected_history_idx == real_idx:
                        _new_meeting()
                    elif (st.session_state.selected_history_idx is not None
                          and st.session_state.selected_history_idx > real_idx):
                        st.session_state.selected_history_idx -= 1
                    st.session_state.renaming_idx = None; st.rerun()
            if st.session_state.renaming_idx == real_idx:
                new_name = st.text_input("이름 바꾸기", value=h.get("title",""),
                    key=f"hist_name_{real_idx}", label_visibility="collapsed")
                sc, cc = st.columns(2)
                with sc:
                    if st.button("저장", key=f"hist_s_{real_idx}", use_container_width=True):
                        st.session_state.meeting_history[real_idx]["title"] = (
                            new_name.strip() or h.get("title","회의"))
                        save_shared_history(st.session_state.meeting_history)
                        st.session_state.renaming_idx = None; st.rerun()
                with cc:
                    if st.button("취소", key=f"hist_cancel_{real_idx}", use_container_width=True):
                        st.session_state.renaming_idx = None; st.rerun()

# ── 메인 레이아웃 ──────────────────────────────────────────
if st.session_state.show_result_panel and st.session_state.last_final:
    main_col, result_col = st.columns([3, 2])
else:
    main_col = st.container(); result_col = None

with main_col:
    col_title, col_cfg = st.columns([5, 1])
    with col_title:
        st.markdown("## 🏢 인사총무팀 AI 에이전트")
        api_status = " ✅" if _api_key else " ⚠️ API키 미입력"
        st.caption(f"설정: {_p.split()[-1]} · {_ml}" + api_status)
    with col_cfg:
        st.write("")
        if st.button("⚙️ 설정", use_container_width=True):
            show_settings()

    if st.session_state.get("stop_notice"):
        st.warning(st.session_state.stop_notice)
        st.session_state.stop_notice = ""

    st.divider()
    table_slot = st.empty()
    with table_slot:
        components.html(
            render_room(st.session_state.bot_statuses,
                        st.session_state.bot_results,
                        st.session_state.current_doc),
            height=350,
        )

    # 업로더를 회의 테이블 위로 겹쳐 올림. 드롭존은 section 태그라 div 한정 셀렉터 금지.
    # 안내문(200MB...)은 숨겨 컴팩트하게. 테이블과 어긋나면 margin 첫 값(-214px)만 조정.
    st.markdown("""
    <style>
    [data-testid="stFileUploader"]{
        width:180px !important; margin:-250px auto -62px auto !important;
        position:relative; z-index:50;
    }
    [data-testid="stFileUploader"] label{
        color:#fff5e1 !important; font-weight:700; font-size:10px;
        display:block !important; width:100% !important;
        text-align:center !important; margin-bottom:2px !important;
    }
    [data-testid="stFileUploaderDropzone"]{
        background:rgba(0,0,0,.16) !important;
        border:1.5px dashed rgba(255,245,225,.85) !important;
        border-radius:9px; min-height:0 !important; height:40px !important;
        width:180px !important; max-width:180px !important; margin:0 auto !important;
        padding:2px 6px !important;
        display:flex; justify-content:center; align-items:center;
    }
    [data-testid="stFileUploaderDropzone"]:hover{ border-color:#ffe6ad !important; }
    [data-testid="stFileUploaderDropzoneInstructions"]{ display:none !important; }
    [data-testid="stFileUploaderDropzone"] button,
    [data-testid="stFileUploaderDropzone"] button *{
        background:#fff7ec !important; color:#5a3d22 !important;
        border:none !important; font-weight:700; margin:0 auto !important;
    }
    </style>
    """, unsafe_allow_html=True)

    ufile = st.file_uploader("📄 클릭 또는 드래그하여 업로드",
                             type=["pdf", "docx", "txt", "pptx"], key="doc_upload")
    if ufile:
        if st.session_state.current_doc != ufile.name:
            st.session_state.current_doc = ufile.name; st.rerun()

    st.divider()
    st.markdown("##### 💬 질문 / 검토 요청")
    question = st.text_area("질문", height=80, label_visibility="collapsed",
        placeholder="예: 이 근로계약서 수습기간 조항이 법적으로 문제없나요?")

    btn_c1, btn_c2 = st.columns([3, 1])
    with btn_c1:
        go = st.button("🚀 회의 시작", type="primary", use_container_width=True,
                       disabled=st.session_state.meeting_running)
    with btn_c2:
        if st.button("⏹ 중단", use_container_width=True,
                     disabled=not st.session_state.meeting_running):
            st.session_state.stop_requested = True
            st.session_state.meeting_running = False; st.rerun()

    # 피드백 섹션
    if st.session_state.last_final and not st.session_state.meeting_running:
        with st.expander("💬 최종 결론에 대한 추가 질문 / 피드백"):
            feedback_q = st.text_area("피드백 내용", key="feedback_input", height=80,
                placeholder="예: 수습 3개월 조항에 대해 더 자세히 알고 싶어요...")
            if st.button("🔄 피드백 반영 추가 회의", key="feedback_go"):
                if feedback_q.strip():
                    st.session_state.feedback_pending = feedback_q.strip(); st.rerun()
                else:
                    st.warning("추가 질문이나 피드백 내용을 먼저 입력해주세요.")

    if st.session_state.get("feedback_pending"):
        go = True

# 최종 결론 패널
if result_col is not None:
    with result_col:
        st.markdown("### 🎯 최종 결론")
        if st.button("✕ 닫기", key="close_result"):
            st.session_state.show_result_panel = False; st.rerun()
        st.markdown("---")
        st.markdown(st.session_state.last_final)

# ── 헬퍼 ──────────────────────────────────────────────────
def upd(statuses_update, results_update=None):
    st.session_state.bot_statuses.update(statuses_update)
    if results_update:
        st.session_state.bot_results.update(results_update)
    with table_slot:
        components.html(
            render_room(st.session_state.bot_statuses,
                        st.session_state.bot_results,
                        st.session_state.current_doc),
            height=350,
        )

def check_stop():
    if st.session_state.stop_requested:
        notice = "⏹ 회의가 중단되었습니다."
        if st.session_state.bot_results:
            from datetime import datetime
            _q = st.session_state.get("current_question", "")
            _pt = "[중단] " + (_q[:25] if _q else st.session_state.current_doc or "회의")
            partial_results = dict(st.session_state.bot_results)
            partial_statuses = dict(st.session_state.bot_statuses)
            partial_summary = build_partial_summary(partial_results)
            entry = {
                "time":     datetime.now().strftime("%Y-%m-%d %H:%M"),
                "title":    _pt,
                "doc":      st.session_state.current_doc,
                "question": _q,
                "ops1": {}, "fc1": "",
                "ops2":     partial_results,
                "fc2": "", "final": "",
                "statuses": partial_statuses,
                "results":  partial_results,
                "partial_summary": partial_summary,
                "partial":  True,
            }
            st.session_state.selected_history_idx = append_shared_history(entry)
            st.session_state.last_final = partial_summary
            st.session_state.show_result_panel = bool(partial_summary)
            notice += " 완료된 검토 내용은 오른쪽 결과 패널과 최근 회의 기록에서 확인할 수 있습니다."
        else:
            notice += " 아직 저장할 완료 검토 내용은 없습니다."
        st.session_state.meeting_running = False
        st.session_state.stop_requested = False
        st.session_state.stop_notice = notice
        st.rerun()

# ── 회의 진행 ──────────────────────────────────────────────
if go:
    if not _api_key:
        st.error("⚙️ 설정에서 API 키를 입력하고 저장해주세요.")
        st.stop()

    _fb = st.session_state.get("feedback_pending", "")
    st.session_state.feedback_pending = ""
    effective_question = question.strip()
    if _fb:
        prior_ctx = ""
        if st.session_state.last_final:
            prior_ctx = f"\n\n[이전 회의 결론 참고]\n{st.session_state.last_final[:600]}"
        effective_question = _fb + prior_ctx

    if not ufile and not effective_question.strip():
        st.error("문서를 업로드하거나 질문을 입력해주세요.")
        st.stop()

    st.session_state.bot_statuses    = {}
    st.session_state.bot_results     = {}
    st.session_state.stop_requested  = False
    st.session_state.meeting_running = True
    st.session_state.show_result_panel = False
    st.session_state.current_question = effective_question

    doc_text = extract_text(ufile) if ufile else ""
    content  = build_content(doc_text, effective_question, st.session_state.case_notes)
    bp       = st.session_state.bot_prompts

    st.divider()
    st.markdown("### 📋 회의 진행")

    active_bots = BOTS[:]
    skipped = []
    with st.spinner("🔍 안건 분석 — 관련 봇 선별 중..."):
        try:
            screen_res = call_bot(_pkey, _api_key, _model, SCREEN_SYS, content[:600], 150)
            m = re.search(r"\{[^{}]+\}", screen_res, re.DOTALL)
            relevant    = json.loads(m.group()) if m else {}
            active_bots = [b for b in BOTS if relevant.get(b["name"], True)]
            skipped     = [b["name"] for b in BOTS if not relevant.get(b["name"], True)]
        except Exception:
            pass
    if skipped:
        st.info("💤 제외된 봇: **" + ", ".join(skipped) + "**")
        upd({n: "skip" for n in skipped})
    n_active = len(active_bots)

    st.markdown(f"**1라운드 — 전문 봇 {n_active}명 동시 발언**")
    upd({b["name"]: "active" for b in active_bots})
    ops1 = {}
    with st.spinner(f"🤖 {n_active}명 동시 분석 중..."):
        def _call1(b):
            return b["name"], call_bot(_pkey, _api_key, _model, bp.get(b["name"], b["system"]), content, 600)
        with ThreadPoolExecutor(max_workers=max(1, n_active)) as ex:
            for name, res in ex.map(_call1, active_bots):
                ops1[name] = res
    upd({b["name"]: "done" for b in active_bots}, ops1)
    for b in active_bots:
        with st.expander(f"🤖 {b['name']} 1차 의견", expanded=False):
            st.markdown(ops1[b["name"]])

    check_stop()

    st.markdown("**팩트체크 1차**")
    fc1_in = content + "\n\n[1차 의견]\n" + "\n\n".join(f"{n}:\n{o}" for n, o in ops1.items())
    upd({"팩트체크봇": "active"})
    with st.status("🤖 팩트체크봇 1차 검증 중...", expanded=False) as stat:
        fc1 = call_bot(_pkey, _api_key, _model, bp.get("팩트체크봇", FACTCHECK_SYS), fc1_in, 900)
        st.markdown(fc1)
        stat.update(label="🤖 팩트체크봇 1차 ✅", state="complete", expanded=False)
    upd({"팩트체크봇": "done"}, {"팩트체크봇_1차": fc1})

    check_stop()

    st.markdown(f"**2라운드 — 팩트체크 반영 ({n_active}명 병렬)**")
    upd({b["name"]: "active" for b in active_bots})
    ops2 = {}
    with st.spinner(f"🤖 {n_active}명 재검토 중..."):
        def _call2(b):
            r2_in = (content + f"\n\n[팩트체크 결과]\n{fc1}"
                     + "\n\n[내 1차 의견]\n" + ops1.get(b["name"], "")
                     + "\n\n팩트체크 반영해 의견을 보완하거나 수정해주세요.")
            return b["name"], call_bot(_pkey, _api_key, _model, bp.get(b["name"], b["system"]), r2_in, 600)
        with ThreadPoolExecutor(max_workers=max(1, n_active)) as ex:
            for name, res in ex.map(_call2, active_bots):
                ops2[name] = res
    upd({b["name"]: "done" for b in active_bots}, ops2)
    for b in active_bots:
        with st.expander(f"🤖 {b['name']} 2차 의견", expanded=False):
            st.markdown(ops2[b["name"]])

    check_stop()

    st.markdown("**팩트체크 2차 (최종 검증)**")
    fc2_in = (content + "\n\n[2차 의견]\n"
              + "\n\n".join(f"{n}:\n{o}" for n, o in ops2.items())
              + f"\n\n[1차 팩트체크]\n{fc1}")
    upd({"팩트체크봇": "active"})
    with st.status("🤖 팩트체크봇 최종 검증 중...", expanded=False) as stat:
        fc2 = call_bot(_pkey, _api_key, _model, bp.get("팩트체크봇", FACTCHECK_SYS), fc2_in, 900)
        st.markdown(fc2)
        stat.update(label="🤖 팩트체크봇 최종 ✅", state="complete", expanded=False)
    upd({"팩트체크봇": "done"}, {"팩트체크봇_2차": fc2})

    check_stop()

    upd({"실무봇": "active"})
    field_in = (content + "\n\n[전문 봇 최종 의견]\n"
                + "\n".join(f"{n}: {o[:200]}" for n, o in ops2.items())
                + f"\n\n[팩트체크 최종]\n{fc2}"
                + "\n\n현장 실무 관점에서 3줄 이내로 핵심 코멘트만 작성해주세요.")
    with st.spinner("🤖 실무봇 코멘트 작성 중..."):
        field_brief = call_bot(_pkey, _api_key, _model, bp.get("실무봇", FIELD_SYS), field_in, 200)
    upd({"실무봇": "done"}, {"실무봇": field_brief})

    final_in = (content + "\n\n[전문 봇 2차 최종 의견]\n"
                + "\n\n".join(f"{n}:\n{o}" for n, o in ops2.items())
                + f"\n\n[팩트체크 최종 결과]\n{fc2}"
                + f"\n\n[실무 현장 코멘트]\n{field_brief}")
    upd({"최종정리봇": "active"})
    with st.spinner("🤖 최종정리봇 정리 중..."):
        final_res = call_bot(_pkey, _api_key, _model, bp.get("최종정리봇", FINAL_SYS), final_in, 1500)
    upd({"최종정리봇": "done"}, {"최종정리봇": final_res})

    st.session_state.meeting_running   = False
    st.session_state.last_final        = final_res
    st.session_state.show_result_panel = True

    from datetime import datetime
    title = (ufile.name if ufile else effective_question[:30])
    entry = {
        "time":     datetime.now().strftime("%Y-%m-%d %H:%M"),
        "title":    title,
        "doc":      ufile.name if ufile else "",
        "question": effective_question[:100],
        "ops1": ops1, "fc1": fc1,
        "ops2": ops2, "fc2": fc2,
        "final": final_res,
        "statuses": dict(st.session_state.bot_statuses),
        "results":  dict(st.session_state.bot_results),
    }
    st.session_state.selected_history_idx = append_shared_history(entry)

    # 자동 메모리
    try:
        mem_in = f"[회의 주제] {title}\n\n[최종 결론 요약]\n{final_res[:800]}"
        memory_note = call_bot(_pkey, _api_key, _model, MEMORY_SYS, mem_in, 120)
        if memory_note and not memory_note.startswith("오류"):
            sep = "\n\n" if st.session_state.case_notes.strip() else ""
            st.session_state.case_notes += sep + f"• {memory_note}"
    except Exception:
        pass

    st.success("✅ 회의 완료! 오른쪽에 최종 결론이 표시됩니다.")
    st.rerun()
# end of app
